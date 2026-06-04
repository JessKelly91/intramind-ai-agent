"""Document Ingestion Workflow using LangGraph.

This workflow handles document ingestion with advanced file parsing,
sophisticated text chunking, and robust error handling.
"""

import logging
import os
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langgraph.graph import END, StateGraph

from config import settings
from models.state import IngestionWorkflowState
from tools import get_api_client
from utils.pii import get_default_redactor

logger = logging.getLogger(__name__)

# Supported file types
SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".txt": "text",
    ".md": "text",
    ".markdown": "text",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".gif": "image",
    ".bmp": "image",
}

# Maximum file size (100MB by default)
MAX_FILE_SIZE_MB = 100
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024


async def validate_document(state: IngestionWorkflowState) -> IngestionWorkflowState:
    """Validate document before processing.
    
    Checks:
    - File exists and is readable
    - File size is within limits
    - File format is supported
    - Required state fields are present
    """
    logger.info("Node: validate_document")

    file_path = state.get("file_path")
    
    # Check file path is provided
    if not file_path:
        error_msg = "No file path provided"
        logger.error(error_msg)
        return {**state, "error": error_msg, "next_step": "handle_error"}

    # Check file exists
    path = Path(file_path)
    if not path.exists():
        error_msg = f"File does not exist: {file_path}"
        logger.error(error_msg)
        return {**state, "error": error_msg, "next_step": "handle_error"}

    # Check file is readable
    if not os.access(file_path, os.R_OK):
        error_msg = f"File is not readable: {file_path}"
        logger.error(error_msg)
        return {**state, "error": error_msg, "next_step": "handle_error"}

    # Check file size
    file_size = path.stat().st_size
    if file_size > MAX_FILE_SIZE_BYTES:
        error_msg = f"File too large: {file_size / (1024*1024):.2f}MB (max: {MAX_FILE_SIZE_MB}MB)"
        logger.error(error_msg)
        return {**state, "error": error_msg, "next_step": "handle_error"}

    if file_size == 0:
        error_msg = "File is empty"
        logger.error(error_msg)
        return {**state, "error": error_msg, "next_step": "handle_error"}

    # Check file format is supported
    file_extension = path.suffix.lower()
    if file_extension not in SUPPORTED_EXTENSIONS:
        error_msg = f"Unsupported file type: {file_extension}. Supported: {', '.join(SUPPORTED_EXTENSIONS.keys())}"
        logger.error(error_msg)
        return {**state, "error": error_msg, "next_step": "handle_error"}

    file_type = SUPPORTED_EXTENSIONS[file_extension]
    
    # Check collection name
    if not state.get("collection_name"):
        error_msg = "No collection name provided"
        logger.error(error_msg)
        return {**state, "error": error_msg, "next_step": "handle_error"}

    logger.info(f"Validation passed - File: {path.name}, Type: {file_type}, Size: {file_size / 1024:.2f}KB")

    return {
        **state,
        "current_step": "validate_document",
        "next_step": "extract_content",
        "document_type": file_type,
        "document_metadata": {
            **state.get("document_metadata", {}),
            "filename": path.name,
            "file_extension": file_extension,
            "file_size_bytes": file_size,
            "file_type": file_type,
        },
    }


async def extract_content(state: IngestionWorkflowState) -> IngestionWorkflowState:
    """Extract text content from document.
    
    Supports:
    - PDF files (pypdf)
    - Word documents (python-docx)
    - PowerPoint presentations (python-pptx)
    - Text files (with encoding detection)
    - Images (basic support, returns metadata)
    """
    logger.info("Node: extract_content")

    file_path = state["file_path"]
    file_type = state.get("document_type", "text")
    path = Path(file_path)

    try:
        if file_type == "pdf":
            content, metadata = _extract_pdf(path)
        elif file_type == "docx":
            content, metadata = _extract_docx(path)
        elif file_type == "pptx":
            content, metadata = _extract_pptx(path)
        elif file_type == "text":
            content, metadata = _extract_text(path)
        elif file_type == "image":
            content, metadata = _extract_image(path)
        else:
            error_msg = f"Unsupported file type for extraction: {file_type}"
            logger.error(error_msg)
            return {**state, "error": error_msg, "next_step": "handle_error"}

        if not content or len(content.strip()) == 0:
            error_msg = "No content extracted from file"
            logger.error(error_msg)
            return {**state, "error": error_msg, "next_step": "handle_error"}

        logger.info(f"Extracted {len(content)} characters from {path.name}")

        return {
            **state,
            "current_step": "extract_content",
            "extracted_content": content,
            "next_step": "redact_pii",
            "document_metadata": {
                **state.get("document_metadata", {}),
                **metadata,
                "extraction_timestamp": datetime.now(timezone.utc).isoformat(),
                "content_length": len(content),
            },
        }

    except Exception as e:
        error_msg = f"Content extraction failed: {str(e)}"
        logger.error(error_msg, exc_info=True)
        return {**state, "error": error_msg, "next_step": "handle_error"}


async def redact_pii(state: IngestionWorkflowState) -> IngestionWorkflowState:
    """Redact PII in extracted content before chunking and storage.

    Step 3 of the Free RAI Stack. Replaces detected PII with stable
    type-tagged tokens so the embedder still sees grammatical text but no
    raw PII ever lands in Weaviate. Findings (type + offsets, no values)
    round-trip via document_metadata for auditability.

    If `settings.enable_pii_redaction` is False or Presidio is unavailable,
    this node is a pass-through and routes straight to chunk_content.
    """
    logger.info("Node: redact_pii")

    content = state.get("extracted_content") or ""
    base_metadata: dict[str, Any] = state.get("document_metadata", {})

    if not settings.enable_pii_redaction:
        logger.info("PII redaction disabled by settings - skipping")
        return {
            **state,
            "current_step": "redact_pii",
            "next_step": "chunk_content",
            "document_metadata": {
                **base_metadata,
                "pii_redaction_applied": False,
                "pii_findings": base_metadata.get("pii_findings", []),
            },
        }

    try:
        redactor = get_default_redactor(
            entities=settings.pii_entities,
            score_threshold=settings.pii_score_threshold,
        )

        if not redactor.available:
            logger.warning(
                "PII redactor unavailable (Presidio/spaCy not installed) - "
                "passing content through unredacted"
            )
            return {
                **state,
                "current_step": "redact_pii",
                "next_step": "chunk_content",
                "document_metadata": {
                    **base_metadata,
                    "pii_redaction_applied": False,
                    "pii_redaction_skipped_reason": "redactor_unavailable",
                    "pii_findings": base_metadata.get("pii_findings", []),
                },
            }

        doc_id = base_metadata.get("filename") or state.get("file_path")
        redacted, findings = redactor.redact(content, doc_id=doc_id)

        existing_findings = base_metadata.get("pii_findings", []) or []
        all_findings = existing_findings + findings

        logger.info(
            "PII redaction complete: %d finding(s), content length %d -> %d",
            len(findings),
            len(content),
            len(redacted),
        )

        return {
            **state,
            "current_step": "redact_pii",
            "extracted_content": redacted,
            "next_step": "chunk_content",
            "document_metadata": {
                **base_metadata,
                "pii_redaction_applied": True,
                "pii_redaction_timestamp": datetime.now(timezone.utc).isoformat(),
                "pii_findings": all_findings,
                "pii_findings_count": len(all_findings),
            },
        }

    except Exception as e:
        error_msg = f"PII redaction failed: {str(e)}"
        logger.error(error_msg, exc_info=True)
        return {**state, "error": error_msg, "next_step": "handle_error"}


def _extract_pdf(path: Path) -> tuple[str, dict[str, Any]]:
    """Extract text from PDF file."""
    from pypdf import PdfReader

    reader = PdfReader(path)
    
    # Extract metadata
    metadata = {
        "page_count": len(reader.pages),
        "pdf_version": reader.pdf_header,
    }
    
    # Extract PDF metadata if available
    if reader.metadata:
        metadata.update({
            "author": reader.metadata.get("/Author", ""),
            "title": reader.metadata.get("/Title", ""),
            "subject": reader.metadata.get("/Subject", ""),
            "creator": reader.metadata.get("/Creator", ""),
        })

    # Extract text from all pages
    text_parts = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text.strip():
            text_parts.append(f"--- Page {i} ---\n{text}")

    content = "\n\n".join(text_parts)
    
    return content, metadata


def _extract_docx(path: Path) -> tuple[str, dict[str, Any]]:
    """Extract text from Word document."""
    from docx import Document

    doc = Document(path)
    
    # Extract metadata
    metadata = {
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
    }
    
    # Extract core properties if available
    if doc.core_properties:
        metadata.update({
            "author": doc.core_properties.author or "",
            "title": doc.core_properties.title or "",
            "subject": doc.core_properties.subject or "",
            "created": doc.core_properties.created.isoformat() if doc.core_properties.created else "",
            "modified": doc.core_properties.modified.isoformat() if doc.core_properties.modified else "",
        })

    # Extract text from paragraphs
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    # Extract text from tables
    for table in doc.tables:
        table_text = []
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                table_text.append(row_text)
        if table_text:
            text_parts.append("\n[TABLE]\n" + "\n".join(table_text) + "\n[/TABLE]\n")

    content = "\n\n".join(text_parts)
    
    return content, metadata


def _extract_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    """Extract text from PowerPoint presentation."""
    from pptx import Presentation

    prs = Presentation(path)
    
    # Extract metadata
    metadata = {
        "slide_count": len(prs.slides),
    }
    
    # Extract core properties if available
    if prs.core_properties:
        metadata.update({
            "author": prs.core_properties.author or "",
            "title": prs.core_properties.title or "",
            "subject": prs.core_properties.subject or "",
            "created": prs.core_properties.created.isoformat() if prs.core_properties.created else "",
            "modified": prs.core_properties.modified.isoformat() if prs.core_properties.modified else "",
        })

    # Extract text from slides
    text_parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = [f"--- Slide {i} ---"]
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        if len(slide_text) > 1:  # More than just the header
            text_parts.append("\n".join(slide_text))

    content = "\n\n".join(text_parts)
    
    return content, metadata


def _extract_text(path: Path) -> tuple[str, dict[str, Any]]:
    """Extract text from plain text file with encoding detection."""
    import chardet

    # Read file as bytes first
    with open(path, "rb") as f:
        raw_data = f.read()

    # Detect encoding
    detected = chardet.detect(raw_data)
    encoding = detected["encoding"] or "utf-8"
    confidence = detected["confidence"]

    logger.info(f"Detected encoding: {encoding} (confidence: {confidence:.2f})")

    # Decode with detected encoding
    try:
        content = raw_data.decode(encoding)
    except UnicodeDecodeError:
        # Fallback to utf-8 with error handling
        logger.warning(f"Failed to decode with {encoding}, falling back to utf-8")
        content = raw_data.decode("utf-8", errors="replace")
        encoding = "utf-8 (fallback)"

    metadata = {
        "encoding": encoding,
        "encoding_confidence": confidence,
        "line_count": content.count("\n") + 1,
    }

    return content, metadata


def _extract_image(path: Path) -> tuple[str, dict[str, Any]]:
    """Extract metadata from image file.
    
    Note: OCR not implemented by default. Returns image info as text.
    For OCR support, uncomment the pytesseract code below.
    """
    from PIL import Image

    img = Image.open(path)
    
    metadata = {
        "image_format": img.format,
        "image_mode": img.mode,
        "image_size": img.size,
        "image_width": img.width,
        "image_height": img.height,
    }

    # Basic content description
    content = f"Image: {path.name}\nFormat: {img.format}\nSize: {img.width}x{img.height}\nMode: {img.mode}"

    # Optional: OCR with pytesseract (requires tesseract installation)
    # Uncomment to enable:
    # try:
    #     import pytesseract
    #     ocr_text = pytesseract.image_to_string(img)
    #     if ocr_text.strip():
    #         content += f"\n\nExtracted Text (OCR):\n{ocr_text}"
    #         metadata["ocr_enabled"] = True
    # except Exception as e:
    #     logger.warning(f"OCR extraction failed: {e}")
    #     metadata["ocr_enabled"] = False

    return content, metadata


async def chunk_content(state: IngestionWorkflowState) -> IngestionWorkflowState:
    """Chunk content using sophisticated text splitting.
    
    Uses LangChain's RecursiveCharacterTextSplitter which:
    - Preserves semantic boundaries (paragraphs, sentences)
    - Handles overlap intelligently
    - Maintains context across chunks
    """
    logger.info("Node: chunk_content")

    content = state["extracted_content"]
    chunk_size = state.get("chunk_size", 1000)
    chunk_overlap = state.get("chunk_overlap", 200)

    try:
        # Create text splitter
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""],
            keep_separator=True,
        )

        # Split text into chunks
        text_chunks = text_splitter.split_text(content)

        # Create chunk objects with metadata
        chunks = []
        for i, chunk_text in enumerate(text_chunks):
            chunks.append({
                "content": chunk_text,
                "index": i,
                "char_count": len(chunk_text),
            })

        logger.info(f"Created {len(chunks)} chunks (size: {chunk_size}, overlap: {chunk_overlap})")

        return {
            **state,
            "current_step": "chunk_content",
            "chunks": chunks,
            "next_step": "store_chunks",
            "document_metadata": {
                **state.get("document_metadata", {}),
                "total_chunks": len(chunks),
                "chunk_size": chunk_size,
                "chunk_overlap": chunk_overlap,
            },
        }

    except Exception as e:
        error_msg = f"Content chunking failed: {str(e)}"
        logger.error(error_msg, exc_info=True)
        return {**state, "error": error_msg, "next_step": "handle_error"}


async def store_chunks(state: IngestionWorkflowState) -> IngestionWorkflowState:
    """Store chunks in vector database with rich metadata.

    Step 6 of the Free RAI Stack: every chunk now carries a standardized
    provenance / governance metadata block:

        source, ingested_by, ingested_at, consent_basis, retention_class,
        pii_findings, pii_redaction_applied

    See ``docs/cards/DATASETS.md`` for the full schema and allowed values.
    """
    logger.info("Node: store_chunks")

    chunks = state["chunks"]
    collection = state["collection_name"]
    client = get_api_client()

    inserted_ids = []
    failed_chunks = []

    try:
        # Prepare documents for batch insertion
        documents = []
        base_metadata = state.get("document_metadata", {}) or {}
        ingested_at = datetime.now(timezone.utc).isoformat()

        # Standardized governance/lineage fields. We don't override values the
        # caller already supplied (e.g. an explicit consent_basis), but we do
        # guarantee every chunk has *something* in each key.
        governance_defaults = {
            "source": base_metadata.get("source")
            or base_metadata.get("filename")
            or state.get("file_path", "unknown"),
            "ingested_by": base_metadata.get("ingested_by", "system"),
            "ingested_at": base_metadata.get("ingested_at", ingested_at),
            "consent_basis": base_metadata.get("consent_basis", "internal"),
            "retention_class": base_metadata.get("retention_class", "standard"),
            "pii_findings": base_metadata.get("pii_findings", []),
            "pii_redaction_applied": base_metadata.get(
                "pii_redaction_applied", False
            ),
        }

        for chunk in chunks:
            doc = {
                "content": chunk["content"],
                "metadata": {
                    **base_metadata,
                    **governance_defaults,
                    "chunk_index": chunk["index"],
                    "total_chunks": len(chunks),
                    "chunk_char_count": chunk["char_count"],
                    "ingestion_timestamp": ingested_at,
                },
            }
            documents.append(doc)

        logger.info(f"Inserting {len(documents)} documents into collection '{collection}'")

        # Batch insert all documents
        result = await client.insert_documents_batch(
            collection_name=collection, 
            documents=documents
        )

        inserted_ids = result.get("ids", [])
        
        if len(inserted_ids) != len(documents):
            logger.warning(f"Partial insertion: {len(inserted_ids)}/{len(documents)} documents inserted")
            failed_chunks = [i for i in range(len(documents)) if i >= len(inserted_ids)]

        logger.info(f"Successfully stored {len(inserted_ids)} chunks in '{collection}'")

        return {
            **state,
            "current_step": "store_chunks",
            "inserted_ids": inserted_ids,
            "workflow_complete": True,
            "document_metadata": {
                **state.get("document_metadata", {}),
                "inserted_count": len(inserted_ids),
                "failed_count": len(failed_chunks),
                "storage_timestamp": datetime.now(timezone.utc).isoformat(),
            },
        }

    except Exception as e:
        error_msg = f"Storage failed: {str(e)}"
        logger.error(error_msg, exc_info=True)
        return {**state, "error": error_msg, "next_step": "handle_error"}


async def handle_error(state: IngestionWorkflowState) -> IngestionWorkflowState:
    """Handle workflow errors."""
    logger.error(f"Ingestion workflow error: {state.get('error')}")
    
    return {
        **state,
        "current_step": "handle_error",
        "workflow_complete": True,
    }


def route_after_validate(state: IngestionWorkflowState) -> str:
    """Route after validation - to error handler or extraction."""
    if state.get("error"):
        return "handle_error"
    return "extract_content"


def route_after_extract(state: IngestionWorkflowState) -> str:
    """Route after extraction - to error handler or PII redaction."""
    if state.get("error"):
        return "handle_error"
    return "redact_pii"


def route_after_redact(state: IngestionWorkflowState) -> str:
    """Route after PII redaction - to error handler or chunking."""
    if state.get("error"):
        return "handle_error"
    return "chunk_content"


def route_after_chunk(state: IngestionWorkflowState) -> str:
    """Route after chunking - to error handler or storage."""
    if state.get("error"):
        return "handle_error"
    return "store_chunks"


def route_after_store(state: IngestionWorkflowState) -> str:
    """Route after storage - to error handler or end."""
    if state.get("error"):
        return "handle_error"
    return "end"


def create_ingestion_workflow() -> StateGraph:
    """Create document ingestion workflow with conditional error routing.

    Workflow flow:
    1. validate_document -> Check file validity
       - If error -> handle_error
       - If success -> extract_content
    2. extract_content -> Extract text from file
       - If error -> handle_error
       - If success -> redact_pii
    3. redact_pii -> Detect + tokenize PII (Step 3: Presidio)
       - If error -> handle_error
       - If success -> chunk_content
    4. chunk_content -> Split into semantic chunks
       - If error -> handle_error
       - If success -> store_chunks
    5. store_chunks -> Store in vector database
       - If error -> handle_error
       - If success -> END
    6. handle_error -> Error handling (routes to END)
    """
    workflow = StateGraph(IngestionWorkflowState)

    # Add nodes
    workflow.add_node("validate_document", validate_document)
    workflow.add_node("extract_content", extract_content)
    workflow.add_node("redact_pii", redact_pii)
    workflow.add_node("chunk_content", chunk_content)
    workflow.add_node("store_chunks", store_chunks)
    workflow.add_node("handle_error", handle_error)

    # Set entry point
    workflow.set_entry_point("validate_document")

    # Add conditional edges that check for errors
    workflow.add_conditional_edges(
        "validate_document",
        route_after_validate,
        {
            "extract_content": "extract_content",
            "handle_error": "handle_error",
        }
    )

    workflow.add_conditional_edges(
        "extract_content",
        route_after_extract,
        {
            "redact_pii": "redact_pii",
            "handle_error": "handle_error",
        }
    )

    workflow.add_conditional_edges(
        "redact_pii",
        route_after_redact,
        {
            "chunk_content": "chunk_content",
            "handle_error": "handle_error",
        }
    )

    workflow.add_conditional_edges(
        "chunk_content",
        route_after_chunk,
        {
            "store_chunks": "store_chunks",
            "handle_error": "handle_error",
        }
    )

    workflow.add_conditional_edges(
        "store_chunks",
        route_after_store,
        {
            "end": END,
            "handle_error": "handle_error",
        }
    )

    # Error handler always ends the workflow
    workflow.add_edge("handle_error", END)

    return workflow.compile()


# Create and export the workflow
ingestion_workflow = create_ingestion_workflow()

